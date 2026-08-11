#!/usr/bin/env python3
"""
Build: Autodesk Platform Services (APS) on AWS — Autodesk University 2026 booth deck.
Modern dark theme, AWS orange accent, booth-readable text, custom architecture shapes.
Output: APS-on-AWS-AU2026.pptx (11 slides)
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.oxml.ns import qn
import os

# ---- Palette ----
BG        = RGBColor(0x0D, 0x0D, 0x0F)   # near-black background
PANEL     = RGBColor(0x1A, 0x1B, 0x20)   # dark panel
PANEL2    = RGBColor(0x24, 0x26, 0x2D)   # lighter panel
ORANGE    = RGBColor(0xFF, 0x99, 0x00)   # AWS orange accent
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
OFFWHITE  = RGBColor(0xE6, 0xE8, 0xEC)
GREY      = RGBColor(0x9A, 0x9F, 0xA8)
BLACK     = RGBColor(0x00, 0x00, 0x00)
BLUE      = RGBColor(0x3B, 0x8E, 0xEA)   # AWS-ish blue for compute
GREEN     = RGBColor(0x3F, 0xB9, 0x50)   # sustainability green
PURPLE    = RGBColor(0x8B, 0x5C, 0xF6)   # identity

# 16:9 widescreen
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def solid(shape, color, line_color=None, line_w=None):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
        shape.line.width = line_w or Pt(1)
    shape.shadow.inherit = False


def bg(slide, color=BG):
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    solid(r, color)
    slide.shapes._spTree.remove(r._element)
    slide.shapes._spTree.insert(2, r._element)
    return r


def accent_bar(slide, top=Inches(1.28), left=Inches(0.7), width=Inches(2.2)):
    b = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, Pt(6))
    solid(b, ORANGE)
    return b


def txt(slide, l, t, w, h, text, size, color=WHITE, bold=False, align=PP_ALIGN.LEFT,
        font="Segoe UI", anchor=MSO_ANCHOR.TOP, italic=False, spacing=None):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    if spacing:
        p.line_spacing = spacing
    r = p.add_run(); r.text = text
    f = r.font
    f.size = Pt(size); f.bold = bold; f.italic = italic
    f.color.rgb = color; f.name = font
    return tb


def bullets(slide, l, t, w, h, items, size=18, color=OFFWHITE, gap=8,
            bullet_color=ORANGE, bold_lead=False):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        p.line_spacing = 1.05
        rd = p.add_run(); rd.text = "▸  "
        rd.font.size = Pt(size); rd.font.color.rgb = bullet_color; rd.font.bold = True
        rd.font.name = "Segoe UI"
        r = p.add_run(); r.text = it
        r.font.size = Pt(size); r.font.color.rgb = color; r.font.name = "Segoe UI"
        r.font.bold = bold_lead
    return tb


def panel(slide, l, t, w, h, color=PANEL, line=None, lw=None, radius=True):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, l, t, w, h)
    solid(shp, color, line, lw)
    return shp


def kicker(slide, text):
    txt(slide, Inches(0.7), Inches(0.55), Inches(9), Inches(0.4),
        text.upper(), 15, ORANGE, bold=True)


def title(slide, text, size=34):
    txt(slide, Inches(0.7), Inches(0.82), Inches(12), Inches(0.9), text, size, WHITE, bold=True)
    accent_bar(slide)


def centered_text_in(shape, lines, sizes, colors, bolds, italics=None, gap=4, spacing=1.0):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Pt(4); tf.margin_right = Pt(4)
    tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    if italics is None:
        italics = [False] * len(lines)
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(gap)
        p.line_spacing = spacing
        r = p.add_run(); r.text = ln
        r.font.size = Pt(sizes[i]); r.font.bold = bolds[i]
        r.font.italic = italics[i]
        r.font.color.rgb = colors[i]; r.font.name = "Segoe UI"


def arrow(slide, l, t, w, h, color=ORANGE):
    a = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, l, t, w, h)
    solid(a, color)
    return a


def page_no(slide, n):
    txt(slide, Inches(12.4), Inches(7.02), Inches(0.7), Inches(0.35),
        f"{n:02d}", 11, GREY, align=PP_ALIGN.RIGHT)


def footer(slide, text="Autodesk Platform Services on AWS  ·  Autodesk University 2026"):
    txt(slide, Inches(0.7), Inches(7.02), Inches(10), Inches(0.35), text, 10, GREY)


# =====================================================================
# Slide 1 — Title
# =====================================================================
s = prs.slides.add_slide(BLANK); bg(s)
# subtle side accent block
sideL = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.18), SH)
solid(sideL, ORANGE)
txt(s, Inches(0.9), Inches(0.7), Inches(6), Inches(0.4),
    "AUTODESK UNIVERSITY 2026", 15, ORANGE, bold=True)
txt(s, Inches(0.9), Inches(2.55), Inches(11.6), Inches(2.0),
    "Autodesk Platform Services on AWS", 54, WHITE, bold=True)
accent_bar(s, top=Inches(4.35), left=Inches(0.95), width=Inches(3.0))
txt(s, Inches(0.95), Inches(4.65), Inches(11.2), Inches(1.0),
    "The design pipeline behind millions of models", 26, OFFWHITE)
txt(s, Inches(0.95), Inches(6.7), Inches(11), Inches(0.4),
    "Powered by AWS  ·  Durable storage  ·  Elastic compute  ·  Global delivery", 14, GREY)

# =====================================================================
# Slide 2 — The story
# =====================================================================
s = prs.slides.add_slide(BLANK); bg(s)
kicker(s, "The Story")
title(s, "From a raw design file to a live 3D model in your browser")
p = panel(s, Inches(0.7), Inches(2.1), Inches(11.93), Inches(2.4), PANEL)
tf = p.text_frame; tf.word_wrap = True
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
tf.margin_left = Inches(0.5); tf.margin_right = Inches(0.5)
para = tf.paragraphs[0]; para.line_spacing = 1.25
r = para.add_run()
r.text = ("A real engineering model, running live in a browser — no CAD software. "
          "It started as a raw design file: uploaded to Autodesk Platform Services, "
          "stored in the cloud, translated to a lightweight web format, and streamed "
          "right here.")
r.font.size = Pt(24); r.font.color.rgb = OFFWHITE; r.font.name = "Segoe UI"
# highlight strip
hp = panel(s, Inches(0.7), Inches(4.85), Inches(11.93), Inches(1.5), PANEL2)
centered_text_in(hp,
    ["Durable storage   ·   Compute that scales on demand   ·   Global low-latency delivery",
     "— all on AWS"],
    [22, 20], [WHITE, ORANGE], [True, True], gap=8)
footer(s); page_no(s, 2)

# =====================================================================
# Slide 3 — What you're seeing (live demo)
# =====================================================================
s = prs.slides.add_slide(BLANK); bg(s)
kicker(s, "Live Demo")
title(s, "What you're seeing right now")
# left: device mock
dev = panel(s, Inches(0.7), Inches(2.1), Inches(5.2), Inches(4.5), PANEL)
scr = panel(s, Inches(1.0), Inches(2.45), Inches(4.6), Inches(3.5), RGBColor(0x10,0x12,0x18),
            line=ORANGE, lw=Pt(1.25))
centered_text_in(scr,
    ["◱", "3D CAD model", "in a plain web browser"],
    [60, 22, 15], [ORANGE, WHITE, GREY], [False, True, False], gap=6)
txt(s, Inches(1.0), Inches(6.05), Inches(4.6), Inches(0.4),
    "Same view on a phone — no plugin, no install", 13, GREY, align=PP_ALIGN.CENTER)
# right: wow bullets as cards
cards = [
    ("No workstation", "No high-end CAD rig — it runs in the cloud."),
    ("Any device", "Laptop, tablet or phone. Just a browser."),
    ("Instant", "Open a link and the model streams in seconds."),
]
cx = Inches(6.25); cw = Inches(6.4); ch = Inches(1.32); cy = 2.1
for h, b in cards:
    cp = panel(s, cx, Inches(cy), cw, ch, PANEL)
    strip = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx, Inches(cy), Inches(0.12), ch)
    solid(strip, ORANGE)
    txt(s, cx + Inches(0.35), Inches(cy) + Inches(0.16), Inches(5.9), Inches(0.4),
        h, 20, WHITE, bold=True)
    txt(s, cx + Inches(0.35), Inches(cy) + Inches(0.62), Inches(5.9), Inches(0.6),
        b, 15, OFFWHITE)
    cy += 1.52
footer(s); page_no(s, 3)

# =====================================================================
# Slide 4 — APS pipeline (5 steps L->R)
# =====================================================================
s = prs.slides.add_slide(BLANK); bg(s)
kicker(s, "How It Works")
title(s, "The APS pipeline — five steps, left to right")
steps = [
    ("1  Upload", ".stl / .step\ndesign file", "Durable storage\nlike Amazon S3", ORANGE),
    ("2  Authenticate", "APS OAuth 2.0\nsecure access", "Enterprise identity\n& access control", PURPLE),
    ("3  Object Storage", "APS OSS\ncloud file store", "Durable object\nstorage (S3)", ORANGE),
    ("4  Model Derivative", "Translate to\nSVF2 web format", "Scalable compute\non demand", BLUE),
    ("5  Stream", "APS Viewer\nin the browser", "Global delivery\nlike CloudFront", GREEN),
]
n = len(steps)
gap = Inches(0.18)
total_w = Inches(11.93)
box_w = Emu(int((total_w - gap * (n - 1)) / n))
x = Inches(0.7); topy = Inches(2.35); box_h = Inches(2.35)
for i, (hd, mid, aws, col) in enumerate(steps):
    bx = Emu(int(x) + i * (int(box_w) + int(gap)))
    card = panel(s, bx, topy, box_w, box_h, PANEL)
    hdr = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bx, topy, box_w, Inches(0.62))
    solid(hdr, col)
    centered_text_in(hdr, [hd], [16], [BLACK if col in (ORANGE, GREEN) else WHITE], [True])
    body = s.shapes.add_textbox(bx, topy + Inches(0.75), box_w, Inches(1.5))
    btf = body.text_frame; btf.word_wrap = True
    btf.margin_left = Pt(6); btf.margin_right = Pt(6)
    for j, ln in enumerate(mid.split("\n")):
        pp = btf.paragraphs[0] if j == 0 else btf.add_paragraph()
        pp.alignment = PP_ALIGN.CENTER; pp.line_spacing = 1.0
        rr = pp.add_run(); rr.text = ln
        rr.font.size = Pt(14); rr.font.color.rgb = OFFWHITE; rr.font.name = "Segoe UI"
        rr.font.bold = True
    # arrow between
    if i < n - 1:
        ax = Emu(int(bx) + int(box_w) - int(Inches(0.02)))
        arrow(s, ax, topy + Inches(0.9), Inches(0.22), Inches(0.4), ORANGE)
# AWS analog row label + cards
txt(s, Inches(0.7), Inches(5.0), Inches(6), Inches(0.4),
    "The AWS building blocks underneath", 16, ORANGE, bold=True)
for i, (hd, mid, aws, col) in enumerate(steps):
    bx = Emu(int(x) + i * (int(box_w) + int(gap)))
    a = panel(s, bx, Inches(5.45), box_w, Inches(1.1), PANEL2)
    centered_text_in(a, aws.split("\n"), [13, 13], [WHITE, GREY], [True, False], gap=2)
footer(s); page_no(s, 4)

# =====================================================================
# Slide 5 — Architecture diagram
# =====================================================================
s = prs.slides.add_slide(BLANK); bg(s)
kicker(s, "Architecture")
title(s, "How the demo is built on AWS")

def arch_box(l, t, w, h, lines, sizes, colors, bolds, fill=PANEL, line=None, lw=None):
    b = panel(s, l, t, w, h, fill, line, lw)
    centered_text_in(b, lines, sizes, colors, bolds, gap=2)
    return b

bh = Inches(1.15)
row_y = Inches(2.45)
# main horizontal chain
chain = [
    (Inches(0.7),  ["Browser", "APS Viewer SDK"], PANEL, ORANGE),
    (Inches(3.15), ["CloudFront", "CDN + TLS"], PANEL, GREEN),
    (Inches(5.6),  ["Cognito", "Sign-in"], PANEL, PURPLE),
    (Inches(8.05), ["ECS Fargate", "APS Node.js app"], PANEL, BLUE),
    (Inches(10.5), ["Autodesk APS APIs", "OSS + Model Derivative"], PANEL, ORANGE),
]
cw = Inches(2.15)
box_refs = []
for l, lines, fill, accent in chain:
    b = panel(s, l, row_y, cw, bh, fill, accent, Pt(1.5))
    centered_text_in(b, lines, [17, 12], [WHITE, GREY], [True, False], gap=2)
    box_refs.append((l, b))
# connecting arrows
for i in range(len(chain) - 1):
    l0 = int(chain[i][0]) + int(cw)
    ax = Emu(l0 - int(Inches(0.02)))
    arrow(s, ax, row_y + Inches(0.42), Inches(0.32), Inches(0.32), ORANGE)
# side box: S3 below Fargate/APS
s3l = Inches(8.05); s3t = Inches(4.55); s3w = Inches(4.6); s3h = Inches(1.25)
s3 = panel(s, s3l, s3t, s3w, s3h, PANEL, ORANGE, Pt(1.5))
centered_text_in(s3, ["Amazon S3", "Model files  ·  pre-signed URLs"],
                 [18, 13], [ORANGE, OFFWHITE], [True, False], gap=3)
# connector from Fargate/APS down to S3 (elbow: simple down arrow)
conn = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(10.15), row_y + bh, Inches(0.32), Inches(1.15))
solid(conn, ORANGE)
# legend
txt(s, Inches(0.7), Inches(6.15), Inches(12), Inches(0.5),
    "Requests flow left→right; model files persist in S3 and are served back to the viewer.",
    15, GREY)
footer(s); page_no(s, 5)

# =====================================================================
# Slide 6 — Persona 1: Developer
# =====================================================================
s = prs.slides.add_slide(BLANK); bg(s)
kicker(s, "Persona 1  ·  Developer")
title(s, "Ship it to AWS in plain English")
bullets(s, Inches(0.7), Inches(2.1), Inches(6.0), Inches(3.5), [
    "Kiro is Amazon's agentic IDE, powered by Bedrock and Claude.",
    "Describe what you want in natural language.",
    "Kiro turns it into infrastructure-as-code.",
    "Deploys ECS Fargate + Amazon S3 + Cognito + CloudFront.",
    "From idea to running app — no hand-written templates.",
], size=19, gap=12)
# code-style callout box
cb = panel(s, Inches(6.95), Inches(2.1), Inches(5.68), Inches(4.4), RGBColor(0x0A,0x0B,0x0E),
           line=ORANGE, lw=Pt(1.25))
# title bar for the callout
tbar = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.95), Inches(2.1), Inches(5.68), Inches(0.5))
solid(tbar, PANEL2)
txt(s, Inches(7.2), Inches(2.17), Inches(5), Inches(0.4),
    "▸ natural-language prompt to Kiro", 13, ORANGE, bold=True, font="Consolas")
ctf = cb.text_frame; ctf.word_wrap = True
ctf.margin_left = Inches(0.35); ctf.margin_right = Inches(0.3); ctf.margin_top = Inches(0.85)
ctf.vertical_anchor = MSO_ANCHOR.TOP
para = ctf.paragraphs[0]; para.line_spacing = 1.25
r = para.add_run()
r.text = ("\"Deploy an APS viewer web app on AWS: a Node.js service on ECS Fargate "
          "behind CloudFront, an S3 bucket for uploaded model files, and Cognito "
          "for sign-in. Wire it to the Autodesk Model Derivative API.\"")
r.font.size = Pt(18); r.font.color.rgb = RGBColor(0x8F,0xE3,0x88); r.font.name = "Consolas"
footer(s); page_no(s, 6)

# =====================================================================
# Slide 7 — Persona 2: Business Analyst
# =====================================================================
s = prs.slides.add_slide(BLANK); bg(s)
kicker(s, "Persona 2  ·  Business Analyst")
title(s, "Design activity becomes business insight")
bullets(s, Inches(0.7), Inches(2.1), Inches(6.0), Inches(3.5), [
    "Amazon QuickSight dashboards over your APS activity.",
    "No CAD software needed — just a browser.",
    "Turn raw pipeline events into decisions.",
    "Embed the same analytics inside your own apps.",
], size=19, gap=13)
# metric tiles
tiles = [
    ("Models translated", "This week"),
    ("Avg translate time", "Per model"),
    ("Storage used", "GB in S3 / OSS"),
    ("Cost", "Per project"),
]
tx = 6.95; ty = 2.1; tw = Inches(2.72); th = Inches(1.55); gx = 0.24; gy = 0.28
for i, (h, b) in enumerate(tiles):
    col = i % 2; rowi = i // 2
    l = Inches(tx + col * (2.72 + gx)); t = Inches(ty + rowi * (1.55 + gy))
    tile = panel(s, l, t, tw, th, PANEL)
    strip = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, tw, Inches(0.1))
    solid(strip, ORANGE)
    centered_text_in(tile, [h, b], [18, 13], [WHITE, GREY], [True, False], gap=6)
footer(s); page_no(s, 7)

# =====================================================================
# Slide 8 — Sustainability + MCP
# =====================================================================
s = prs.slides.add_slide(BLANK); bg(s)
kicker(s, "Sustainability  ·  MCP")
title(s, "Ask your designs about their carbon footprint")
# flow: Autodesk Sustainability API -> MCP server -> Bedrock / QuickSight
flow = [
    (["Autodesk", "Sustainability API"], GREEN),
    (["MCP Server", "Model Context Protocol"], ORANGE),
    (["Bedrock + QuickSight", "Natural-language answers"], BLUE),
]
fx = Inches(0.7); fy = Inches(2.35); fw = Inches(3.55); fh = Inches(1.5); fgap = Inches(0.55)
for i, (lines, accent) in enumerate(flow):
    bx = Emu(int(fx) + i * (int(fw) + int(fgap)))
    b = panel(s, bx, fy, fw, fh, PANEL, accent, Pt(1.5))
    centered_text_in(b, lines, [19, 13], [WHITE, GREY], [True, False], gap=4)
    if i < len(flow) - 1:
        ax = Emu(int(bx) + int(fw) + int(Inches(0.06)))
        arrow(s, ax, fy + Inches(0.55), Inches(0.42), Inches(0.4), ORANGE)
# explanation panel
ep = panel(s, Inches(0.7), Inches(4.35), Inches(11.93), Inches(1.55), PANEL)
centered_text_in(ep,
    ["Understand the carbon footprint of materials and designs — in plain language.",
     "Ask: \"Which material choice lowers this part's footprint the most?\""],
    [20, 18], [OFFWHITE, ORANGE], [True, False], italics=[False, True], gap=8)
txt(s, Inches(0.7), Inches(6.2), Inches(12), Inches(0.5),
    "Already working today: the Autodesk Product Help MCP server connects to Kiro.",
    15, GREEN, bold=True)
footer(s); page_no(s, 8)

# =====================================================================
# Slide 9 — Why AWS (4 boxes)
# =====================================================================
s = prs.slides.add_slide(BLANK); bg(s)
kicker(s, "Why AWS")
title(s, "Four reasons this runs so well on AWS")
boxes = [
    ("Durable & secure storage", "Model files kept safe and highly available.", ORANGE),
    ("Compute that scales to zero", "Pay for translation only when it runs.", BLUE),
    ("Global low-latency delivery", "Stream models fast, anywhere in the world.", GREEN),
    ("Enterprise identity & governance", "Sign-in, access control, and audit built in.", PURPLE),
]
bw = Inches(5.9); bh2 = Inches(2.0); bx0 = 0.7; by0 = 2.2; gxx = 0.33; gyy = 0.35
for i, (h, b, col) in enumerate(boxes):
    col_i = i % 2; row_i = i // 2
    l = Inches(bx0 + col_i * (5.9 + gxx)); t = Inches(by0 + row_i * (2.0 + gyy))
    card = panel(s, l, t, bw, bh2, PANEL)
    strip = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, Inches(0.14), bh2)
    solid(strip, col)
    txt(s, l + Inches(0.45), t + Inches(0.3), bw - Inches(0.8), Inches(0.7),
        h, 23, WHITE, bold=True)
    txt(s, l + Inches(0.45), t + Inches(1.15), bw - Inches(0.8), Inches(0.7),
        b, 17, OFFWHITE)
footer(s); page_no(s, 9)

# =====================================================================
# Slide 10 — Three booth demos
# =====================================================================
s = prs.slides.add_slide(BLANK); bg(s)
kicker(s, "At The Booth")
title(s, "Three things you can try today")
demos = [
    ("a", "Live APS Viewer app", "Open a real 3D model in a plain browser — on your own phone.", ORANGE),
    ("b", "Autodesk MCP + Kiro / Bedrock", "Ask questions and deploy to AWS in natural language.", BLUE),
    ("c", "QuickSight analytics + sustainability", "See pipeline metrics and design carbon footprint.", GREEN),
]
dy = 2.25
for tag, h, b, col in demos:
    l = Inches(0.7); t = Inches(dy); w = Inches(11.93); hh = Inches(1.4)
    card = panel(s, l, t, w, hh, PANEL)
    badge = s.shapes.add_shape(MSO_SHAPE.OVAL, l + Inches(0.35), t + Inches(0.38), Inches(0.62), Inches(0.62))
    solid(badge, col)
    centered_text_in(badge, [tag], [26], [BLACK if col in (ORANGE, GREEN) else WHITE], [True])
    txt(s, l + Inches(1.3), t + Inches(0.24), Inches(10.2), Inches(0.55),
        h, 24, WHITE, bold=True)
    txt(s, l + Inches(1.3), t + Inches(0.82), Inches(10.2), Inches(0.5),
        b, 16, OFFWHITE)
    dy += 1.62
footer(s); page_no(s, 10)

# =====================================================================
# Slide 11 — Closing / CTA
# =====================================================================
s = prs.slides.add_slide(BLANK); bg(s)
sideL = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.18), SH)
solid(sideL, ORANGE)
txt(s, Inches(0.9), Inches(1.15), Inches(8), Inches(0.4),
    "AUTODESK PLATFORM SERVICES ON AWS", 15, ORANGE, bold=True)
txt(s, Inches(0.9), Inches(2.0), Inches(8.4), Inches(2.6),
    "Designs stored, translated, and streamed to any device — worldwide.",
    38, WHITE, bold=True)
accent_bar(s, top=Inches(4.5), left=Inches(0.95), width=Inches(3.0))
txt(s, Inches(0.95), Inches(4.85), Inches(8.2), Inches(0.8),
    "Come build with us.", 30, ORANGE, bold=True)
# QR placeholder box
qr = panel(s, Inches(9.9), Inches(2.0), Inches(2.7), Inches(2.7), PANEL, ORANGE, Pt(1.5))
centered_text_in(qr, ["▢ QR", "scan for the", "demo repo"],
                 [40, 15, 15], [ORANGE, GREY, GREY], [True, False, False], gap=4)
txt(s, Inches(9.9), Inches(4.85), Inches(2.7), Inches(0.5),
    "github.com/…/aps-on-aws", 13, GREY, align=PP_ALIGN.CENTER)
footer(s)

# ---- save ----
OUT = "/Users/batramc/.kiro/crew/workspace/aps-aws-demo/deck/APS-on-AWS-AU2026.pptx"
prs.save(OUT)
print("SAVED", OUT)
print("SLIDES", len(prs.slides._sldIdLst))
