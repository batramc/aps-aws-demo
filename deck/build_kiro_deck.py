#!/usr/bin/env python3
"""
Build: Built on Kiro — Autodesk Platform Services on AWS (Autodesk University 2026).
Kiro-centric: what Kiro does and how it built + deploys the APS demo. No comparisons.
Same dark AWS-orange theme as the APS deck. Output: Built-on-Kiro-AU2026.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BG=RGBColor(0x0D,0x0D,0x0F); PANEL=RGBColor(0x1A,0x1B,0x20); PANEL2=RGBColor(0x24,0x26,0x2D)
ORANGE=RGBColor(0xFF,0x99,0x00); WHITE=RGBColor(0xFF,0xFF,0xFF); OFFWHITE=RGBColor(0xE6,0xE8,0xEC)
GREY=RGBColor(0x9A,0x9F,0xA8); BLACK=RGBColor(0,0,0); BLUE=RGBColor(0x3B,0x8E,0xEA)
GREEN=RGBColor(0x3F,0xB9,0x50); PURPLE=RGBColor(0x8B,0x5C,0xF6); TEAL=RGBColor(0x2A,0xB7,0xCA)

prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
SW,SH=prs.slide_width,prs.slide_height; BLANK=prs.slide_layouts[6]

def solid(sh,c,lc=None,lw=None):
    sh.fill.solid(); sh.fill.fore_color.rgb=c
    if lc is None: sh.line.fill.background()
    else: sh.line.color.rgb=lc; sh.line.width=lw or Pt(1)
    sh.shadow.inherit=False

def bg(s,c=BG):
    r=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,SW,SH); solid(r,c)
    s.shapes._spTree.remove(r._element); s.shapes._spTree.insert(2,r._element)

def txt(s,l,t,w,h,text,size,color=WHITE,bold=False,align=PP_ALIGN.LEFT,font="Segoe UI",
        anchor=MSO_ANCHOR.TOP,italic=False,spacing=None):
    tb=s.shapes.add_textbox(l,t,w,h); tf=tb.text_frame; tf.word_wrap=True
    tf.vertical_anchor=anchor
    tf.margin_left=0;tf.margin_right=0;tf.margin_top=0;tf.margin_bottom=0
    p=tf.paragraphs[0]; p.alignment=align
    if spacing: p.line_spacing=spacing
    r=p.add_run(); r.text=text; f=r.font
    f.size=Pt(size); f.bold=bold; f.italic=italic; f.color.rgb=color; f.name=font
    return tb

def bullets(s,l,t,w,h,items,size=18,color=OFFWHITE,gap=8,bc=ORANGE):
    tb=s.shapes.add_textbox(l,t,w,h); tf=tb.text_frame; tf.word_wrap=True
    tf.margin_left=0;tf.margin_right=0;tf.margin_top=0;tf.margin_bottom=0
    for i,it in enumerate(items):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.space_after=Pt(gap); p.line_spacing=1.08
        rd=p.add_run(); rd.text="▸  "; rd.font.size=Pt(size); rd.font.color.rgb=bc
        rd.font.bold=True; rd.font.name="Segoe UI"
        r=p.add_run(); r.text=it; r.font.size=Pt(size); r.font.color.rgb=color; r.font.name="Segoe UI"
    return tb

def panel(s,l,t,w,h,c=PANEL,line=None,lw=None,radius=True):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,l,t,w,h)
    solid(sh,c,line,lw); return sh

def center(shape,lines,sizes,colors,bolds,italics=None,gap=4,spacing=1.0):
    tf=shape.text_frame; tf.word_wrap=True; tf.vertical_anchor=MSO_ANCHOR.MIDDLE
    tf.margin_left=Pt(6);tf.margin_right=Pt(6);tf.margin_top=Pt(2);tf.margin_bottom=Pt(2)
    if italics is None: italics=[False]*len(lines)
    for i,ln in enumerate(lines):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment=PP_ALIGN.CENTER; p.space_after=Pt(gap); p.line_spacing=spacing
        r=p.add_run(); r.text=ln; r.font.size=Pt(sizes[i]); r.font.bold=bolds[i]
        r.font.italic=italics[i]; r.font.color.rgb=colors[i]; r.font.name="Segoe UI"

def accent(s,top=Inches(1.28),left=Inches(0.7),width=Inches(2.2)):
    solid(s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,left,top,width,Pt(6)),ORANGE)

def kicker(s,text): txt(s,Inches(0.7),Inches(0.55),Inches(11),Inches(0.4),text.upper(),15,ORANGE,bold=True)
def title(s,text,size=32): txt(s,Inches(0.7),Inches(0.82),Inches(12),Inches(0.9),text,size,WHITE,bold=True); accent(s)
def arrow(s,l,t,w,h,c=ORANGE): solid(s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,l,t,w,h),c)
def foot(s,n): txt(s,Inches(0.7),Inches(7.02),Inches(10),Inches(0.35),"Built on Kiro  ·  APS on AWS  ·  Autodesk University 2026",10,GREY); txt(s,Inches(12.4),Inches(7.02),Inches(0.7),Inches(0.35),f"{n:02d}",11,GREY,align=PP_ALIGN.RIGHT)

# ---- S1 Title ----
s=prs.slides.add_slide(BLANK); bg(s)
solid(s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,Inches(0.18),SH),ORANGE)
txt(s,Inches(0.9),Inches(0.7),Inches(9),Inches(0.4),"AUTODESK UNIVERSITY 2026  ·  DEVELOPER PERSONA",15,ORANGE,bold=True)
txt(s,Inches(0.9),Inches(2.5),Inches(11.8),Inches(1.6),"Built on Kiro",54,WHITE,bold=True)
accent(s,top=Inches(4.2),left=Inches(0.95),width=Inches(3.0))
txt(s,Inches(0.95),Inches(4.5),Inches(11.4),Inches(1.0),"From idea to a live Autodesk Platform Services app on AWS — with Kiro.",26,OFFWHITE)
txt(s,Inches(0.95),Inches(6.7),Inches(11),Inches(0.4),"Kiro — Amazon's agentic IDE, powered by Amazon Bedrock",14,GREY)

# ---- S2 What is Kiro ----
s=prs.slides.add_slide(BLANK); bg(s); kicker(s,"What Is Kiro"); title(s,"Amazon's agentic IDE — describe it, plan it, build it, ship it")
steps=[("Describe","Say what you want\nin plain English",ORANGE),
       ("Plan","Kiro writes the spec\nbefore any code",BLUE),
       ("Build","It implements,\ntask by task",PURPLE),
       ("Ship","Deploy to AWS\nfrom the IDE",GREEN)]
n=len(steps); fx=Inches(0.7); fy=Inches(2.5); gap=Inches(0.3)
fw=Emu(int((Inches(11.93)-gap*(n-1))/n)); fh=Inches(2.0)
for i,(hd,mid,col) in enumerate(steps):
    bx=Emu(int(fx)+i*(int(fw)+int(gap)))
    b=panel(s,bx,fy,fw,fh,PANEL,col,Pt(1.5))
    center(b,[hd]+mid.split("\n"),[22,14,14],[col,OFFWHITE,GREY],[True,False,False],gap=4)
    if i<n-1: arrow(s,Emu(int(bx)+int(fw)-int(Inches(0.02))),fy+Inches(0.8),Inches(0.3),Inches(0.4),ORANGE)
p=panel(s,Inches(0.7),Inches(5.05),Inches(11.93),Inches(1.4),PANEL2)
center(p,["One agentic IDE takes you from a sentence to a running app on AWS —",
          "powered by Amazon Bedrock, so your code and context stay in your AWS boundary."],
       [20,18],[WHITE,ORANGE],[True,False],gap=8,spacing=1.1)
foot(s,2)

# ---- S3 Spec-driven (headliner) ----
s=prs.slides.add_slide(BLANK); bg(s); kicker(s,"Signature Capability"); title(s,"Spec-driven development — the plan is visible before code")
flow=[("requirements.md","Structured\nrequirements (EARS)",ORANGE),
      ("design.md","The technical\napproach",BLUE),
      ("tasks.md","A checklist Kiro\nexecutes one by one",GREEN)]
fx=Inches(0.7); fy=Inches(2.5); fw=Inches(3.55); fh=Inches(2.1); gap=Inches(0.55)
for i,(hd,mid,col) in enumerate(flow):
    bx=Emu(int(fx)+i*(int(fw)+int(gap)))
    b=panel(s,bx,fy,fw,fh,PANEL,col,Pt(1.5))
    center(b,[hd]+mid.split("\n"),[20,14,14],[col,OFFWHITE,GREY],[True,False,False],gap=4)
    if i<len(flow)-1: arrow(s,Emu(int(bx)+int(fw)+int(Inches(0.06))),fy+Inches(0.85),Inches(0.42),Inches(0.4),ORANGE)
p=panel(s,Inches(0.7),Inches(5.05),Inches(11.93),Inches(1.4),PANEL2)
center(p,["You approve the plan, THEN Kiro builds.",
          "Planning and guardrails are visible and reviewable — no black box, no surprises."],
       [21,16],[WHITE,OFFWHITE],[True,False],gap=8,spacing=1.1)
foot(s,3)

# ---- S4 Agent Hooks + Steering (two feature cards) ----
s=prs.slides.add_slide(BLANK); bg(s); kicker(s,"Automation & Consistency"); title(s,"Agent Hooks and Steering keep work automatic and on-standard")
L=panel(s,Inches(0.7),Inches(2.2),Inches(5.85),Inches(4.2),PANEL,ORANGE,Pt(1.5))
txt(s,Inches(1.05),Inches(2.5),Inches(5.2),Inches(0.5),"Agent Hooks",24,ORANGE,bold=True)
bullets(s,Inches(1.05),Inches(3.25),Inches(5.2),Inches(3.0),[
    "Event-driven automation inside the IDE.",
    "\"On file save, run the tests.\"",
    "\"On API change, update the docs.\"",
    "Repetitive steps happen for you.",
],size=17,gap=12,bc=ORANGE)
R=panel(s,Inches(6.78),Inches(2.2),Inches(5.85),Inches(4.2),PANEL,PURPLE,Pt(1.5))
txt(s,Inches(7.13),Inches(2.5),Inches(5.2),Inches(0.5),"Steering Files",24,PURPLE,bold=True)
bullets(s,Inches(7.13),Inches(3.25),Inches(5.2),Inches(3.0),[
    "Persistent project rules the agent always follows.",
    "Coding standards and conventions.",
    "Architecture constraints and guardrails.",
    "Every generation stays on-standard.",
],size=17,gap=12,bc=PURPLE)
foot(s,4)

# ---- S5 MCP ----
s=prs.slides.add_slide(BLANK); bg(s); kicker(s,"Connect Your Tools"); title(s,"MCP support — connect Autodesk's own data to Kiro")
flow=[(["Autodesk","Product Help MCP"],GREEN),(["Kiro","agentic IDE"],ORANGE),(["Amazon Bedrock","Claude"],BLUE)]
fx=Inches(0.7); fy=Inches(2.45); fw=Inches(3.55); fh=Inches(1.6); fgap=Inches(0.55)
for i,(lines,col) in enumerate(flow):
    bx=Emu(int(fx)+i*(int(fw)+int(fgap)))
    b=panel(s,bx,fy,fw,fh,PANEL,col,Pt(1.5))
    center(b,lines,[20,14],[col,GREY],[True,False],gap=4)
    if i<len(flow)-1: arrow(s,Emu(int(bx)+int(fw)+int(Inches(0.06))),fy+Inches(0.6),Inches(0.42),Inches(0.4),ORANGE)
p=panel(s,Inches(0.7),Inches(4.4),Inches(11.93),Inches(1.5),PANEL)
center(p,["Model Context Protocol lets Kiro talk to Autodesk's own services and APIs —",
          "ask product questions, pull docs, and act on them, right inside the IDE."],
       [20,18],[OFFWHITE,ORANGE],[True,False],gap=8,spacing=1.1)
txt(s,Inches(0.7),Inches(6.2),Inches(12),Inches(0.5),"Already working today: the Autodesk Product Help MCP server connected to Kiro.",15,GREEN,bold=True)
foot(s,5)

# ---- S6 Ship to AWS in plain English ----
s=prs.slides.add_slide(BLANK); bg(s); kicker(s,"Ship It"); title(s,"Deploy the APS app to AWS — in plain English")
bullets(s,Inches(0.7),Inches(2.2),Inches(6.0),Inches(3.6),[
    "Describe the infrastructure you want.",
    "Kiro turns it into infrastructure-as-code.",
    "Deploys App Runner / ECS Fargate + Amazon S3.",
    "Adds Amazon Cognito sign-in and CloudFront.",
    "Live HTTPS URL in minutes — no hand-written YAML.",
],size=19,gap=12)
cb=panel(s,Inches(6.95),Inches(2.2),Inches(5.68),Inches(4.2),RGBColor(0x0A,0x0B,0x0E),ORANGE,Pt(1.25))
tbar=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(6.95),Inches(2.2),Inches(5.68),Inches(0.5)); solid(tbar,PANEL2)
txt(s,Inches(7.2),Inches(2.27),Inches(5),Inches(0.4),"▸ natural-language prompt to Kiro",13,ORANGE,bold=True,font="Consolas")
ctf=cb.text_frame; ctf.word_wrap=True; ctf.margin_left=Inches(0.35); ctf.margin_right=Inches(0.3); ctf.margin_top=Inches(0.85)
pp=ctf.paragraphs[0]; pp.line_spacing=1.25
r=pp.add_run(); r.text=("\"Containerize this Autodesk Platform Services viewer app and deploy it to "
    "AWS App Runner. Put my APS Client ID and Secret in Secrets Manager, create the "
    "IAM roles, and give me the live HTTPS URL.\"")
r.font.size=Pt(18); r.font.color.rgb=RGBColor(0x8F,0xE3,0x88); r.font.name="Consolas"
foot(s,6)

# ---- S7 Why on AWS ----
s=prs.slides.add_slide(BLANK); bg(s); kicker(s,"Why It Matters"); title(s,"Kiro is AWS-native — structured, visible, governed")
boxes=[("AWS-native","Runs on Amazon Bedrock; code and context stay in your AWS boundary.",ORANGE),
       ("Structured","Spec-driven — a reviewable plan before any code.",BLUE),
       ("Automated","Agent Hooks handle the repetitive steps for you.",GREEN),
       ("On-standard","Steering files enforce your rules on every change.",PURPLE)]
bw=Inches(5.9); bh=Inches(2.0); bx0=0.7; by0=2.2; gxx=0.33; gyy=0.35
for i,(h,b,col) in enumerate(boxes):
    ci=i%2; ri=i//2
    l=Inches(bx0+ci*(5.9+gxx)); t=Inches(by0+ri*(2.0+gyy))
    panel(s,l,t,bw,bh,PANEL)
    solid(s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,l,t,Inches(0.14),bh),col)
    txt(s,l+Inches(0.45),t+Inches(0.3),bw-Inches(0.8),Inches(0.7),h,23,WHITE,bold=True)
    txt(s,l+Inches(0.45),t+Inches(1.1),bw-Inches(0.8),Inches(0.8),b,17,OFFWHITE,spacing=1.1)
foot(s,7)

# ---- S8 The whole demo, built on Kiro ----
s=prs.slides.add_slide(BLANK); bg(s); kicker(s,"At The Booth"); title(s,"This entire demo was built on Kiro")
demos=[("a","The APS viewer app","Generated and wired to the Model Derivative API with Kiro.",ORANGE),
       ("b","The AWS deployment","App Runner + S3 + Cognito, described in plain English.",BLUE),
       ("c","The Autodesk MCP integration","Product Help MCP connected to Kiro, live today.",GREEN)]
dy=2.25
for tag,h,b,col in demos:
    l=Inches(0.7); t=Inches(dy); w=Inches(11.93); hh=Inches(1.4)
    panel(s,l,t,w,hh,PANEL)
    badge=s.shapes.add_shape(MSO_SHAPE.OVAL,l+Inches(0.35),t+Inches(0.38),Inches(0.62),Inches(0.62)); solid(badge,col)
    center(badge,[tag],[26],[BLACK if col in (ORANGE,GREEN) else WHITE],[True])
    txt(s,l+Inches(1.3),t+Inches(0.24),Inches(10.2),Inches(0.55),h,24,WHITE,bold=True)
    txt(s,l+Inches(1.3),t+Inches(0.82),Inches(10.2),Inches(0.5),b,16,OFFWHITE)
    dy+=1.62
foot(s,8)

# ---- S9 Close ----
s=prs.slides.add_slide(BLANK); bg(s)
solid(s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,Inches(0.18),SH),ORANGE)
txt(s,Inches(0.9),Inches(1.15),Inches(9),Inches(0.4),"BUILT ON KIRO  ·  ON AWS",15,ORANGE,bold=True)
txt(s,Inches(0.9),Inches(2.1),Inches(11.6),Inches(2.2),"Describe it. Watch the plan. Ship it to AWS.",42,WHITE,bold=True)
accent(s,top=Inches(4.7),left=Inches(0.95),width=Inches(3.0))
txt(s,Inches(0.95),Inches(5.05),Inches(11),Inches(0.8),"Come build with us.",30,ORANGE,bold=True)
foot(s,9)

OUT="/Users/batramc/.kiro/crew/workspace/aps-aws-demo/deck/Built-on-Kiro-AU2026.pptx"
prs.save(OUT); print("SAVED",OUT); print("SLIDES",len(prs.slides._sldIdLst))
